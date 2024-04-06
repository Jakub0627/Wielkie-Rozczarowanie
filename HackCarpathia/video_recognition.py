import cv2
from moviepy.editor import VideoFileClip
import numpy as np
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os


def find_slide_changes_and_save_frames(video_path, threshold=0.01, min_interval=1, averaging_frames=5):
    cap = cv2.VideoCapture(video_path)
    frame_queue = []
    timestamps = []
    last_change = -min_interval
    prev_avg_frame = None
    saved_frames = []

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break
        timestamp = cap.get(cv2.CAP_PROP_POS_MSEC) / 1000.0

        frame_gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
        frame_queue.append(frame_gray)
        if len(frame_queue) > averaging_frames:
            frame_queue.pop(0)

        if len(frame_queue) == averaging_frames:
            avg_frame = np.mean(np.array(frame_queue), axis=0).astype(np.uint8)
            if prev_avg_frame is not None:
                diff = cv2.absdiff(avg_frame, prev_avg_frame)
                non_zero_count = np.count_nonzero(diff)
                total_pixels = diff.size
                change_ratio = non_zero_count / total_pixels

                if change_ratio > threshold and (timestamp - last_change) > min_interval:
                    timestamps.append(timestamp)
                    last_change = timestamp
                    img_path = tempfile.mktemp(suffix='.jpg')
                    cv2.imwrite(img_path, frame)
                    saved_frames.append(img_path)

            prev_avg_frame = avg_frame

    cap.release()
    return timestamps, saved_frames


def create_pptx_with_frames(frame_paths, pptx_path='presentation.pptx'):
    prs = Presentation()
    for img_path in frame_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width)
        os.remove(img_path)

    prs.save(pptx_path)


def cut_video_segments(video_path, timestamps, output_folder):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    clip = VideoFileClip(video_path)
    timestamps.append(clip.duration)

    for i, start_time in enumerate(timestamps[:-1]):
        end_time = timestamps[i + 1]
        segment = clip.subclip(start_time, end_time)
        output_path = os.path.join(output_folder, f'segment_{i}.mp4')
        segment.write_videofile(output_path, codec="libx264")
